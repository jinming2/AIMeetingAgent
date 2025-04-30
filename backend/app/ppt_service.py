import os
from fastapi import HTTPException
import pptx
from openai import OpenAI
import logging
import traceback

# 配置日志
logger = logging.getLogger(__name__)


class PPTService:
    def __init__(self):
        self.openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        if not self.openai_client.api_key:
            raise ValueError("OpenAI API key not found in environment variables")
        logger.info("PPTService initialized with OpenAI client")

    def extract_ppt_content(self, ppt_path: str) -> str:
        """
        从PPT文件中提取文本内容
        """
        try:
            logger.info(f"Attempting to extract content from PPT file: {ppt_path}")
            presentation = pptx.Presentation(ppt_path)
            content = []

            for slide in presentation.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)

                if slide_content:
                    content.append("\n".join(slide_content))

            extracted_content = "\n\n".join(content)
            logger.info(
                f"Successfully extracted content from PPT file. Content length: {len(extracted_content)}"
            )
            return extracted_content
        except Exception as e:
            error_msg = (
                f"Error extracting PPT content: {str(e)}\n{traceback.format_exc()}"
            )
            logger.error(error_msg)
            raise HTTPException(
                status_code=500, detail=f"Error processing PPT file: {str(e)}"
            )

    async def generate_outline(self, ppt_content: str) -> dict:
        """
        使用GPT生成PPT大纲
        """
        try:
            logger.info("Starting to generate outline with GPT")
            prompt = f"""
            你是一个演讲教练，正在帮助演讲者为以下PPT内容准备演讲。

            请生成一个**详细的演讲大纲**，不仅仅列出目录，还要写出每一部分中演讲者可以讲什么，讲的顺序是什么，应该重点强调哪些内容。

            请严格按照如下格式输出：

            1. 【主要主题】  
            - 简要介绍整体目标和听众应当获得的内容
            2. 【子主题1：XXX】
            - 开场引导（如：你可以这样开始…）
            - 主要讲点 1（内容 + 示例/解释）
            - 主要讲点 2（内容 + 示例/解释）
            - 小结提示（例如可以说：“这一部分我们讲了…”）
            3. 【子主题2：XXX】
            - （重复上述结构）

            ……
            N. 【结尾总结】
            - 总结内容要点
            - 强调行动呼吁或目标落点
            - 推荐结语方式

            要求：
            - 内容应来自PPT提取的文本
            - 请避免空泛总结，尽量将提取内容组织成条理清晰的语言结构
            - 可以适当补充合理的过渡语和演讲节奏建议

            以下是PPT内容（按顺序排列）：
            -----------------------------
            {ppt_content}
            """

            logger.info("Sending request to OpenAI API")
            response = self.openai_client.chat.completions.create(
                model="gpt-4o", messages=[{"role": "user", "content": prompt}]
            )

            outline = response.choices[0].message.content
            logger.info("Successfully generated outline")

            return {"outline": outline, "status": "success"}
        except Exception as e:
            error_msg = f"Error generating outline: {str(e)}\n{traceback.format_exc()}"
            logger.error(error_msg)
            raise HTTPException(
                status_code=500, detail=f"Error generating outline: {str(e)}"
            )


# 创建全局实例
ppt_service = PPTService()
