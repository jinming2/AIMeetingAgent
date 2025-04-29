import os
import logging
from pptx import Presentation
import openai

def load_api_key():
    """
    Load OpenAI API key from environment variable.
    """
    key = os.getenv("OPENAI_API_KEY")
    if not key:
        raise ValueError("Missing OPENAI_API_KEY environment variable")
    openai.api_key = key

class PPTAutoSummarizer:
    """
    A utility class to extract content from PPTX files and generate a structured outline using OpenAI.
    """
    def __init__(self, model: str = "gpt-4o-mini"):
        load_api_key()
        self.model = model
        self.logger = logging.getLogger(self.__class__.__name__)
        self.logger.setLevel(logging.INFO)
        if not self.logger.hasHandlers():
            handler = logging.StreamHandler()
            handler.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(message)s"))
            self.logger.addHandler(handler)

    def extract_slides(self, ppt_path: str) -> list[str]:
        """
        Extract text content from each slide in the PPTX file.
        Returns a list of raw slide texts.
        """
        self.logger.info(f"Extracting text from PPT: {ppt_path}")
        presentation = Presentation(ppt_path)
        slides = []
        for slide in presentation.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slide_text = "\n".join(texts)
            if slide_text:
                slides.append(slide_text)
        self.logger.info(f"Extracted {len(slides)} slides with content")
        return slides

    def summarize_ppt(self, ppt_path: str) -> str:
        """
        Extract all slides, combine into one文本,并生成一个结构化大纲（仅大纲，Bullet Point形式）。
        """
        slides = self.extract_slides(ppt_path)
        combined_text = "\n\n".join(slides)

        prompt = (
            "你是演讲者的助手，请根据以下幻灯片内容生成一个结构化大纲，"
            "使用数字分级（1., 2., ...）和短横线子项目符号（-），格式示例如下：\n\n"
            "1. 第一大点\n"
            "   - 第一小点\n"
            "   - 第二小点\n"
            "2. 第二大点\n"
            "   - 第一小点\n"
            "...（如此类推）\n\n"
            "幻灯片内容：\n" + combined_text
        )

        self.logger.info("Generating structured outline for the entire PPT")
        response = openai.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
        )
        outline = response.choices[0].message.content.strip()
        return outline
